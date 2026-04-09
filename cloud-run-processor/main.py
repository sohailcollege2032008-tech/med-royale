"""
Dactoor Question Processor — Cloud Run Service v1.2
- PDF / images  → Gemini Files API  (full fidelity)
- PPTX          → inline base64 (text + embedded images extracted with python-pptx)
- DOCX          → inline base64 text (python-docx)
"""

import base64
import json
import logging
import os
import tempfile
import time

import google.generativeai as genai
from fastapi import FastAPI, File, Header, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware

# ── Logging ────────────────────────────────────────────────────────────────────
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger(__name__)

# ── Config ─────────────────────────────────────────────────────────────────────
GEMINI_API_KEY  = os.environ["GEMINI_API_KEY"]
API_SECRET      = os.environ.get("API_SECRET", "")
ALLOWED_ORIGINS = os.environ.get("ALLOWED_ORIGINS", "*").split(",")
MODEL_NAME      = os.environ.get("GEMINI_MODEL", "gemini-2.0-flash")

genai.configure(api_key=GEMINI_API_KEY)

# ── Safety Settings ────────────────────────────────────────────────────────────
# Set to BLOCK_NONE to allow medical content and avoid false positives
SAFETY_SETTINGS = [
    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
]

# ── FastAPI ────────────────────────────────────────────────────────────────────
app = FastAPI(title="Dactoor Processor", version="1.3.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)

# ── System Prompt ──────────────────────────────────────────────────────────────
SYSTEM_PROMPT = """\
You are a medical exam question extractor. You receive a document \
(PDF, PPTX, DOCX, image, etc.) that contains multiple-choice questions (MCQs).

Your task is to extract ALL questions from the document and return them in \
this EXACT JSON format. Return ONLY valid JSON, no markdown, no explanation.

{
  "title": "<infer a title from the document content>",
  "questions": [
    {
      "id": 1,
      "question": "<the question text in its original language>",
      "question_ar": "<Arabic version if the original is in Arabic, otherwise null>",
      "choices": ["<choice A>", "<choice B>", "<choice C>", "<choice D>"],
      "correct": <0-indexed position of the correct answer>,
      "needs_image": false,
      "image_url": null
    }
  ]
}

RULES:
1. Extract every single MCQ from the document — do not skip any.
2. The "correct" field must be the 0-based index of the correct answer in the choices array.
3. If the correct answer is marked/highlighted/bolded/starred, use that. If no answer is marked, set "correct" to -1.
4. If the question is in Arabic, put it in both "question" and "question_ar". If in English, put in "question" only and set "question_ar" to null.
5. Preserve the original wording of questions and choices exactly as written.
6. If choices are labeled A/B/C/D or 1/2/3/4, remove the labels and just keep the text.
7. Set "needs_image" to true if the question refers to a figure, image, photograph, diagram, \
graph, table, or any visual element that is required to answer correctly — even if the image \
was already provided above. Set to false otherwise.
8. Return ONLY the JSON object. No markdown backticks, no commentary.\
"""

# ── MIME helpers ───────────────────────────────────────────────────────────────
EXT_TO_MIME = {
    "pdf":  "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt":  "application/vnd.ms-powerpoint",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc":  "application/msword",
    "jpg":  "image/jpeg",
    "jpeg": "image/jpeg",
    "png":  "image/png",
    "gif":  "image/gif",
    "webp": "image/webp",
    "bmp":  "image/bmp",
}

# Gemini Files API supports these; everything else goes inline or via extraction
GEMINI_FILES_SUPPORTED = {
    "application/pdf",
    "image/jpeg", "image/png", "image/gif", "image/webp",
    "image/heic", "image/heif", "image/bmp",
}

# Image MIME types Gemini can render inline
INLINE_IMAGE_MIMES = {"image/jpeg", "image/png", "image/gif", "image/webp"}

MAX_INLINE_IMAGES = 30   # cap to avoid huge requests


# ── PPTX extractor ─────────────────────────────────────────────────────────────
def extract_pptx(path: str):
    """
    Returns (text_summary: str, images: list[dict])
    images = [{"mime_type": str, "data": bytes}, ...]
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    text_parts = []
    images = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            # Text
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            # Embedded picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    mime = img.content_type or "image/png"
                    if mime not in INLINE_IMAGE_MIMES:
                        mime = "image/png"
                    images.append({"mime_type": mime, "data": img.blob, "slide": i})
                except Exception as e:
                    log.warning(f"Could not extract image from slide {i}: {e}")

        if slide_texts:
            text_parts.append(f"=== Slide {i} ===\n" + "\n".join(slide_texts))

    if not text_parts and not images:
        raise ValueError("No readable content found in the PPTX file")

    return "\n\n".join(text_parts), images


# ── DOCX extractor ─────────────────────────────────────────────────────────────
def extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt and txt not in lines:
                    lines.append(txt)
    if not lines:
        raise ValueError("No text found in the DOCX file")
    return "\n".join(lines)


# ── JSON helpers ───────────────────────────────────────────────────────────────
def parse_gemini_response(raw: str) -> dict:
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        lines = cleaned.splitlines()
        cleaned = "\n".join(lines[1:])
        if cleaned.rstrip().endswith("```"):
            cleaned = cleaned.rstrip()[:-3]
    return json.loads(cleaned.strip())

# ── Fallback Config ────────────────────────────────────────────────────────────
MODEL_FALLBACK_LIST = [
    "gemini-3.1-flash-lite-preview",
    "gemini-2.5-flash-lite",
    "gemini-2-flash-lite",
    "gemini-3-flash",
    "gemini-2-flash",
    "gemma-4-31b",
    "gemma-4-26b"
]

SAFETY_SETTINGS = [
    {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
]

def generate_with_fallback(contents, system_prompt):
    """
    Attempts generation with fallback sequence.
    Returns (data, model_used, attempts_log)
    """
    attempts_log = []
    
    for model_id in MODEL_FALLBACK_LIST:
        try:
            log.info(f"Attempting with model: {model_id}")
            model = genai.GenerativeModel(
                model_name=model_id,
                system_instruction=system_prompt
            )
            
            response = model.generate_content(contents, safety_settings=SAFETY_SETTINGS)
            
            # Check for safety blocks (finish_reason 4 = RECITING)
            if response.candidates and response.candidates[0].finish_reason == 4:
                reason_msg = "Safety block: Copyright/Recitation"
                log.warning(f"Model {model_id} failed: {reason_msg}")
                attempts_log.append({"model": model_id, "status": "blocked", "reason": reason_msg})
                continue

            # Standard success path
            raw_text = response.text
            data = parse_gemini_response(raw_text)
            return data, model_id, attempts_log

        except Exception as e:
            err_msg = str(e)
            log.warning(f"Model {model_id} error: {err_msg}")
            attempts_log.append({"model": model_id, "status": "error", "reason": err_msg})
            continue

    raise HTTPException(
        status_code=500,
        detail={
            "message": "All AI models failed to process this document.",
            "attempts": attempts_log
        }
    )


# ── Health ─────────────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "version": "1.3.0", "fallback_models": MODEL_FALLBACK_LIST}


# ── Main endpoint ──────────────────────────────────────────────────────────────
@app.post("/process")
async def process_file(
    file: UploadFile = File(...),
    x_api_secret: str | None = Header(None, alias="x-api-secret"),
):
    if API_SECRET and x_api_secret != API_SECRET:
        raise HTTPException(status_code=401, detail="Unauthorized — invalid x-api-secret header")

    content  = await file.read()
    filename = file.filename or "document"
    mime     = file.content_type or "application/octet-stream"
    ext      = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    log.info(f"Received: {filename}  {mime}  {len(content)/1024:.1f} KB")

    if len(content) > 150 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 150 MB)")

    if mime in ("application/octet-stream", "binary/octet-stream"):
        mime = EXT_TO_MIME.get(ext, mime)

    tmp_path      = None
    uploaded_file = None
    
    try:
        suffix = f".{ext}" if ext else ""
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        # Prepare contents list
        contents = []

        # ── Path A: PDF or image → Gemini Files API ────────────────────────────
        if mime in GEMINI_FILES_SUPPORTED:
            log.info(f"Path A (Files API): {mime}")
            uploaded_file = genai.upload_file(tmp_path, mime_type=mime, display_name=filename)

            poll = 0
            while uploaded_file.state.name == "PROCESSING" and poll < 90:
                time.sleep(2)
                uploaded_file = genai.get_file(uploaded_file.name)
                poll += 1

            if uploaded_file.state.name != "ACTIVE":
                raise HTTPException(
                    status_code=500,
                    detail=f"Gemini file processing failed (state={uploaded_file.state.name})"
                )
            contents = [uploaded_file]

        # ── Path B: PPTX → extract text + images ────────────────────────────
        elif ext in ("pptx", "ppt"):
            log.info("Path B (PPTX multimodal)")
            extracted_text, slide_images = extract_pptx(tmp_path)
            contents = [f"Below is the text extracted from the PowerPoint file '{filename}':\n\n{extracted_text}"]
            capped = slide_images[:MAX_INLINE_IMAGES]
            for img in capped:
                contents.append({
                    "mime_type": img["mime_type"],
                    "data": base64.b64encode(img["data"]).decode()
                })

        # ── Path C: DOCX → extract text ────────────────────────
        elif ext in ("docx", "doc"):
            log.info("Path C (DOCX text)")
            extracted_text = extract_docx(tmp_path)
            contents = [f"Below is the text extracted from the Word document '{filename}':\n\n{extracted_text}"]

        else:
            raise HTTPException(
                status_code=415,
                detail=f"Unsupported file type: .{ext}. Supported: PDF, PPTX, DOCX, and images."
            )

        # ── Run with fallback ──────────────────────────────────────────────────
        data, model_used, attempts = generate_with_fallback(contents, SYSTEM_PROMPT)

        # ── Post-process and return ───────────────────────────────────────────
        if not isinstance(data.get("title"), str) or not isinstance(data.get("questions"), list):
            raise HTTPException(status_code=422, detail="AI returned unexpected JSON structure")
        
        data["questions"] = [{**q, "id": i + 1} for i, q in enumerate(data["questions"])]
        data["model_used"] = model_used
        data["is_rollback"] = (model_used != MODEL_FALLBACK_LIST[0])
        data["attempts_log"] = attempts

        log.info(f"✓ {len(data['questions'])} questions extracted via {model_used}")
        return data

    except HTTPException:
        raise
    except Exception as exc:
        log.exception("Unexpected processing error")
        raise HTTPException(status_code=500, detail=str(exc))

    finally:
        if tmp_path:
            try: os.unlink(tmp_path)
            except Exception: pass
        if uploaded_file:
            try: genai.delete_file(uploaded_file.name)
            except Exception: pass
